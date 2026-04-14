from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DIAGRAMS = ROOT / "diagrams"

SLIDES = [
    ("Webmining", "Arquitectura, datos y optimizacion de recorridos"),
    ("Agenda", "Vision, arquitectura, datos Google, optimizacion y roadmap"),
    ("Problema de Negocio", "Planificar recorridos eficientes y factibles"),
    ("Propuesta de Solucion", "Aplicacion Flet + motor de optimizacion + pipeline de datos"),
    ("DFD Nivel 0", "Vista de contexto"),
    ("DFD Nivel 1", "Procesos principales"),
    ("DFD Nivel 2", "Detalle de optimizacion"),
    ("Arquitectura de Contexto", "Interaccion entre sistema y externos"),
    ("Arquitectura de Contenedores", "Separacion por capas"),
    ("Arquitectura de Componentes", "Modulo a modulo"),
    ("Pipeline de Datos Google", "Foco principal de ingenieria de datos"),
    ("Google Places", "Descubrimiento de POIs por categoria"),
    ("Google Distance Matrix", "Construccion de matrices de costo"),
    ("Calidad de Datos", "Normalizacion, deduplicacion y validacion"),
    ("Modelo de Datos (DER)", "Trazabilidad de entidades y resultados"),
    ("Flujo de Optimizacion", "De input a ruta final"),
    ("Resultados y Exportes", "Visualizacion, CSV y PDF"),
    ("Riesgos y Mitigaciones", "Cuotas API, escalabilidad y operacion"),
    ("Roadmap", "Proximas fases tecnicas"),
    ("Cierre", "Impacto y siguientes pasos"),
]

DIAGRAM_FOR_SLIDE = {
    4: DIAGRAMS / "dfd_level_0.png",
    5: DIAGRAMS / "dfd_level_1.png",
    6: DIAGRAMS / "dfd_level_2_optimizacion.png",
    7: DIAGRAMS / "architecture_context.png",
    8: DIAGRAMS / "architecture_containers.png",
    9: DIAGRAMS / "architecture_components.png",
    10: DIAGRAMS / "flow_google_data_pipeline.png",
    14: DIAGRAMS / "der_webmining.png",
    15: DIAGRAMS / "flow_optimizacion.png",
}

BULLETS_FOR_SLIDE = {
    1: [
        "Contexto y objetivos del proyecto",
        "Arquitectura de software y decisiones de diseno",
        "Pipeline de datos Google (enfasis)",
        "Optimizacion, resultados y roadmap",
    ],
    2: [
        "Armado manual de recorridos es costoso y subjetivo",
        "Necesidad de balancear distancia, tiempo y valor turistico",
        "Restricciones horarias vuelven el problema mas realista",
    ],
    3: [
        "UI moderna en Flet (desktop, web y movil)",
        "Motor de optimizacion configurable",
        "Reportes operativos para toma de decisiones",
    ],
    11: [
        "Consulta por area geografica y categorias de interes",
        "Captura de atributos: nombre, tipo, coordenadas",
        "Control de duplicados y consistencia",
    ],
    12: [
        "Consulta de pares origen-destino",
        "Construccion de matriz de distancias y tiempos",
        "Base cuantitativa para optimizacion",
    ],
    13: [
        "Estandarizacion de tipos y horarios",
        "Eliminacion de outliers y registros invalidos",
        "Versionado de datasets por depot",
    ],
    16: [
        "Metricas de distancia, tiempo y puntaje",
        "Mapa de recorrido y trazabilidad por tramo",
        "Exportes CSV/PDF para uso operativo",
    ],
    17: [
        "Cuotas y costos API -> cache + retry/backoff",
        "Escalabilidad combinatoria -> heuristicas y poda",
        "Calidad de datos -> validadores y QA de datasets",
    ],
    18: [
        "Fase 1: robustez y testing",
        "Fase 2: escalabilidad del motor",
        "Fase 3: producto y API multiusuario",
    ],
    19: [
        "El mayor valor tecnico esta en el pipeline de datos Google",
        "La arquitectura modular permite evolucion controlada",
        "Proximo paso: industrializar calidad y performance",
    ],
}

AUDIENCE_PROFILES = {
    "ejecutiva": {
        "output": ROOT / "Webmining_Arquitectura_y_Datos_Google_20min.pptx",
        "bg": (10, 27, 51),
        "title": (245, 249, 255),
        "subtitle": (191, 220, 255),
        "bullet": (236, 244, 255),
        "card": (27, 49, 77),
        "card_border": (94, 138, 188),
        "card_title": (208, 228, 255),
        "footer": "Presentacion ejecutiva para stakeholders y negocio.",
    },
    "docente": {
        "output": ROOT / "Webmining_Presentacion_Docente_20min.pptx",
        "bg": (24, 42, 74),
        "title": (250, 252, 255),
        "subtitle": (204, 225, 255),
        "bullet": (242, 248, 255),
        "card": (35, 66, 110),
        "card_border": (122, 170, 226),
        "card_title": (224, 238, 255),
        "footer": "Version docente: enfoque didactico y trazabilidad de decisiones.",
    },
    "tribunal_tecnico": {
        "output": ROOT / "Webmining_Presentacion_Tribunal_Tecnico_20min.pptx",
        "bg": (14, 22, 31),
        "title": (235, 255, 244),
        "subtitle": (173, 224, 202),
        "bullet": (222, 245, 233),
        "card": (26, 47, 41),
        "card_border": (78, 154, 127),
        "card_title": (202, 240, 223),
        "footer": "Version tribunal tecnico: profundidad en arquitectura, datos y riesgos.",
    },
}


def add_background(slide, color: tuple[int, int, int]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def add_title(slide, title: str, subtitle: str | None, profile: dict[str, object]) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*profile["title"])
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(*profile["subtitle"])


def add_bullets(slide, items: list[str], profile: dict[str, object], x=0.8, y=1.7, w=6.4, h=5.0) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(*profile["bullet"])
        first = False


def add_visual_card(slide, title: str, image_path: Path, profile: dict[str, object], x=7.2, y=1.7, w=5.7, h=5.0) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*profile["card"])
    card.line.color.rgb = RGBColor(*profile["card_border"])

    t = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12), Inches(w - 0.4), Inches(0.4))
    tp = t.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(14)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(*profile["card_title"])

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(x + 0.25), Inches(y + 0.55), Inches(w - 0.5), Inches(h - 0.8))


def add_footer(slide, text: str, profile: dict[str, object]) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(*profile["subtitle"])


def build_deck(name: str, profile: dict[str, object]) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]

    for idx, (title, subtitle) in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)
        add_background(slide, profile["bg"])
        add_title(slide, title, subtitle, profile)

        if idx in BULLETS_FOR_SLIDE:
            add_bullets(slide, BULLETS_FOR_SLIDE[idx], profile)

        if idx in DIAGRAM_FOR_SLIDE:
            add_visual_card(slide, "Vista tecnica", DIAGRAM_FOR_SLIDE[idx], profile)

        add_footer(slide, profile["footer"], profile)

        note = slide.notes_slide.notes_text_frame
        note.text = f"[{name}] Slide {idx + 1}/20 - tiempo sugerido: 45-75 segundos."

    output = profile["output"]
    prs.save(output)
    print(f"Presentacion generada: {output}")


def build_all() -> None:
    for audience, profile in AUDIENCE_PROFILES.items():
        build_deck(audience, profile)


if __name__ == "__main__":
    build_all()
